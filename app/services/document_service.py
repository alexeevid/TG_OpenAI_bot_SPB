# app/services/document_service.py
from __future__ import annotations

import base64
import io
import logging
import mimetypes
import os
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

import fitz  # PyMuPDF
from PIL import Image
from pypdf import PdfReader
from docx import Document as DocxDocument
import openpyxl
from bs4 import BeautifulSoup
from pptx import Presentation

log = logging.getLogger(__name__)


@dataclass
class ExtractResult:
    text: str
    info: str
    warnings: List[str]
    kind: str = "text"          # text / document / table / image / mixed
    description: str = ""       # для image или когда OCR пустой


class DocumentService:
    """
    Единый сервис извлечения текста/содержимого из файлов + OCR + нормализация/сжатие.

    Важно:
    - Для изображений: сначала классификация (документ/таблица vs фото/сцена).
      Если текста нет — возвращаем description, чтобы бот мог описать, что на картинке.
    - Совместимо с минимальным Dockerfile (без apt-get).
    """

    def __init__(self, openai_client, settings):
        self._openai = openai_client
        self._cfg = settings

        self.max_pdf_pages = int(getattr(settings, "max_pdf_pages", 12))
        self.max_zip_files = int(getattr(settings, "max_zip_files", 12))
        self.max_zip_total_mb = int(getattr(settings, "max_zip_total_mb", 25))
        self.max_table_rows = int(getattr(settings, "max_table_rows", 60))
        self.max_table_cols = int(getattr(settings, "max_table_cols", 20))

        self.max_chars_before_compress = int(getattr(settings, "max_chars_before_compress", 18000))
        self.target_chars_after_compress = int(getattr(settings, "target_chars_after_compress", 9000))

        self.max_image_side = int(getattr(settings, "ocr_max_image_side", 1600))

    # ---------------- public ----------------

    def extract_text(self, path: str | Path, *, filename: str = "", mime: Optional[str] = None) -> ExtractResult:
        p = Path(path)
        if not p.exists():
            return ExtractResult("", "file_not_found", ["file_not_found"], kind="text")

        filename = filename or p.name
        mime = self._guess_mime(filename, mime)

        try:
            if mime in ("application/zip", "application/x-zip-compressed") or p.suffix.lower() == ".zip":
                res = self._extract_zip(p)
                return self._postprocess(res)

            if mime.startswith("image/"):
                res = self._extract_image_vision(p)  # <- важное изменение
                return self._postprocess(res)

            if mime == "application/pdf" or p.suffix.lower() == ".pdf":
                res = self._extract_pdf(p)
                return self._postprocess(res)

            ext = p.suffix.lower()

            if ext == ".docx":
                res = self._extract_docx(p)
                return self._postprocess(res)

            if ext in (".xlsx", ".xlsm", ".xltx", ".xltm"):
                res = self._extract_xlsx(p)
                return self._postprocess(res)

            if ext == ".pptx":
                res = self._extract_pptx(p)
                return self._postprocess(res)

            if ext in (".html", ".htm"):
                res = self._extract_html(p)
                return self._postprocess(res)

            if ext in (".csv", ".txt", ".md"):
                res = self._extract_textfile(p)
                return self._postprocess(res)

            res = self._extract_textfile(p)
            return self._postprocess(res)

        except Exception as e:
            log.exception("DocumentService.extract_text failed: %s", e)
            return ExtractResult("", f"error:{type(e).__name__}", [f"exception:{type(e).__name__}"], kind="text")

    # ---------------- postprocess ----------------

    def _postprocess(self, res: ExtractResult) -> ExtractResult:
        """
        Нормализация + сжатие больших текстов.
        ВАЖНО: description для kind=image не трогаем.
        """
        if res.kind == "image" and not res.text.strip():
            # фото/сцена: оставляем description, текст пустой — это нормально
            return res

        text = (res.text or "").strip()
        if not text:
            return res

        text = self._normalize_text(text)
        if len(text) <= self.max_chars_before_compress:
            res.text = text
            return res

        compressed = self._compress_with_llm(text, target_chars=self.target_chars_after_compress)
        if compressed:
            res.warnings.append(f"compressed:{len(text)}->{len(compressed)}")
            res.text = compressed.strip()
            res.info = res.info + " +compressed"
            return res

        cut = text[: self.target_chars_after_compress].rstrip() + "\n...\n[TRUNCATED]"
        res.warnings.append("compressed:fallback_truncate")
        res.text = cut
        res.info = res.info + " +truncated"
        return res

    def _normalize_text(self, text: str) -> str:
        t = text.replace("\r\n", "\n").replace("\r", "\n")
        t = re.sub(r"\n{4,}", "\n\n\n", t)
        lines = [ln.rstrip() for ln in t.split("\n")]
        out: List[str] = []
        prev = None
        dup_count = 0
        for ln in lines:
            if prev is not None and ln == prev and ln.strip():
                dup_count += 1
                if dup_count <= 1:
                    out.append(ln)
                continue
            dup_count = 0
            out.append(ln)
            prev = ln
        t = "\n".join(out).strip()
        t = re.sub(r"[ \t]{3,}", "  ", t)
        return t

    def _compress_with_llm(self, text: str, *, target_chars: int) -> str:
        model = self._text_model_for_compress()
        prompt = (
            "Сожми текст документа, сохранив смысл, структуру и ключевые факты.\n"
            f"- Итоговый объём: примерно до {target_chars} символов.\n"
            "- Сохраняй заголовки/разделы.\n"
            "- Если есть таблицы — оставь их в текстовом виде.\n"
            "- Сохраняй цифры, сроки, критерии, формулировки требований.\n"
            "- Убери повторы и воду.\n"
            "Верни ТОЛЬКО сжатый текст без комментариев.\n"
        )

        if len(text) > 60000:
            head = text[:30000]
            tail = text[-20000:]
            text_in = head + "\n...\n[SKIPPED]\n...\n" + tail
        else:
            text_in = text

        messages = [{"role": "user", "content": prompt + "\n\n---\nТЕКСТ:\n" + text_in + "\n---\n"}]
        try:
            out = self._openai.generate_text(
                model=model,
                messages=messages,
                temperature=0.0,
                max_output_tokens=int(getattr(self._cfg, "openai_max_output_tokens", 1800) or 1800),
                reasoning_effort=getattr(self._cfg, "openai_reasoning_effort", None),
            )
            return (out or "").strip()
        except Exception as e:
            log.warning("LLM compress failed: %s", e)
            return ""

    def _text_model_for_compress(self) -> str:
        m = getattr(self._cfg, "document_compress_model", None)
        if m:
            return str(m)
        return str(getattr(self._cfg, "openai_text_model", "gpt-4o-mini"))

    # ---------------- extractors ----------------

    def _guess_mime(self, filename: str, mime: Optional[str]) -> str:
        if mime:
            return mime
        mt, _ = mimetypes.guess_type(filename or "")
        return mt or "application/octet-stream"

    def _vision_model(self) -> str:
        m = getattr(self._cfg, "openai_vision_model", None)
        if m:
            return str(m)
        return str(getattr(self._cfg, "openai_text_model", "gpt-4o-mini"))

    def _extract_textfile(self, p: Path) -> ExtractResult:
        data = p.read_bytes()
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("cp1251", errors="replace")
        text = text.strip()
        return ExtractResult(text, f"text:{p.suffix.lower()} chars={len(text)}", [], kind="text")

    def _extract_docx(self, p: Path) -> ExtractResult:
        doc = DocxDocument(str(p))
        parts: List[str] = []
        for para in doc.paragraphs:
            t = (para.text or "").strip()
            if t:
                parts.append(t)
        text = "\n".join(parts).strip()
        return ExtractResult(text, f"docx paragraphs={len(doc.paragraphs)} chars={len(text)}", [], kind="document")

    def _extract_xlsx(self, p: Path) -> ExtractResult:
        wb = openpyxl.load_workbook(str(p), data_only=True)
        out: List[str] = []
        for ws in wb.worksheets[:3]:
            out.append(f"## Sheet: {ws.title}")
            max_row = min(ws.max_row or 0, self.max_table_rows)
            max_col = min(ws.max_column or 0, self.max_table_cols)
            for r in range(1, max_row + 1):
                row_vals: List[str] = []
                for c in range(1, max_col + 1):
                    v = ws.cell(r, c).value
                    row_vals.append("" if v is None else str(v).strip())
                if any(x for x in row_vals):
                    out.append(" | ".join(row_vals))
        text = "\n".join(out).strip()
        return ExtractResult(text, f"xlsx sheets={len(wb.worksheets)} chars={len(text)}", [], kind="table")

    def _extract_pptx(self, p: Path) -> ExtractResult:
        pres = Presentation(str(p))
        out: List[str] = []
        for i, slide in enumerate(pres.slides[:30], start=1):
            out.append(f"## Slide {i}")
            for shape in slide.shapes:
                txt = getattr(shape, "text", None)
                if txt:
                    t = str(txt).strip()
                    if t:
                        out.append(t)
        text = "\n".join(out).strip()
        return ExtractResult(text, f"pptx slides={min(len(pres.slides),30)} chars={len(text)}", [], kind="document")

    def _extract_html(self, p: Path) -> ExtractResult:
        data = p.read_bytes()
        try:
            html = data.decode("utf-8")
        except UnicodeDecodeError:
            html = data.decode("cp1251", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n")
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        return ExtractResult(text, f"html chars={len(text)}", [], kind="document")

    def _extract_pdf(self, p: Path) -> ExtractResult:
        try:
            reader = PdfReader(str(p))
            parts: List[str] = []
            for i, page in enumerate(reader.pages[: self.max_pdf_pages]):
                t = ""
                try:
                    t = (page.extract_text() or "").strip()
                except Exception:
                    t = ""
                if t:
                    parts.append(f"## Page {i+1}\n{t}")
            text = "\n\n".join(parts).strip()
            if text:
                return ExtractResult(
                    text,
                    f"pdf:text pages={min(len(reader.pages),self.max_pdf_pages)} chars={len(text)}",
                    [],
                    kind="document",
                )
        except Exception as e:
            log.warning("pypdf extract failed, will try render OCR: %s", e)

        return self._extract_pdf_render_ocr(p)

    def _extract_pdf_render_ocr(self, p: Path) -> ExtractResult:
        doc = fitz.open(str(p))
        n_pages = min(doc.page_count, self.max_pdf_pages)
        warnings: List[str] = []
        if doc.page_count > n_pages:
            warnings.append(f"pdf_pages_limited:{doc.page_count}->{n_pages}")

        out_parts: List[str] = []
        for i in range(n_pages):
            page = doc.load_page(i)
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            # OCR только текст
            text = self._vision_extract(img_bytes)["text"]
            if text.strip():
                out_parts.append(f"## Page {i+1}\n{text.strip()}")

        joined = "\n\n".join(out_parts).strip()
        if not joined:
            warnings.append("pdf_ocr_empty")
        return ExtractResult(joined, f"pdf:render_ocr pages={n_pages} chars={len(joined)}", warnings, kind="document")

    def _extract_zip(self, p: Path) -> ExtractResult:
        warnings: List[str] = []
        total_bytes = p.stat().st_size
        if total_bytes > self.max_zip_total_mb * 1024 * 1024:
            warnings.append(f"zip_too_large:{total_bytes}")

        out_parts: List[str] = []
        count = 0
        total_unpacked = 0

        with zipfile.ZipFile(str(p), "r") as z:
            names = [n for n in z.namelist() if not n.endswith("/")]
            if len(names) > self.max_zip_files:
                warnings.append(f"zip_files_limited:{len(names)}->{self.max_zip_files}")
                names = names[: self.max_zip_files]

            for name in names:
                count += 1
                if ".." in Path(name).parts:
                    warnings.append(f"zip_skip_unsafe:{name}")
                    continue

                try:
                    data = z.read(name)
                except Exception:
                    warnings.append(f"zip_read_failed:{name}")
                    continue

                total_unpacked += len(data)
                if total_unpacked > self.max_zip_total_mb * 1024 * 1024:
                    warnings.append("zip_unpacked_limit_reached")
                    break

                tmp = Path("/tmp") / f"zip_{os.getpid()}_{count}_{Path(name).name}"
                try:
                    tmp.write_bytes(data)
                    child = self.extract_text(tmp, filename=Path(name).name, mime=None)
                    if child.text.strip() or child.description.strip():
                        payload = child.text.strip() or child.description.strip()
                        out_parts.append(f"# File: {name}\n{payload}")
                    else:
                        warnings.append(f"zip_child_empty:{name}")
                finally:
                    try:
                        tmp.unlink(missing_ok=True)
                    except Exception:
                        pass

        text = "\n\n".join(out_parts).strip()
        return ExtractResult(text, f"zip files={count} chars={len(text)}", warnings, kind="mixed")

    # ---------------- Vision for images ----------------

    def _extract_image_vision(self, p: Path) -> ExtractResult:
        img = Image.open(str(p)).convert("RGB")
        w, h = img.size
        scale = min(1.0, self.max_image_side / max(w, h))
        if scale < 1.0:
            img = img.resize((int(w * scale), int(h * scale)))

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        img_bytes = buf.getvalue()

        data = self._vision_extract(img_bytes)

        kind = (data.get("kind") or "image").strip().lower()
        text = (data.get("text") or "").strip()
        desc = (data.get("description") or "").strip()

        if kind not in ("document", "table", "image", "mixed"):
            kind = "image"

        info = f"image:vision kind={kind} chars={len(text)}"
        # если kind=document/table, но текста ноль — всё равно полезно вернуть описание
        return ExtractResult(text=text, info=info, warnings=[], kind=kind, description=desc)

    def _vision_extract(self, img_bytes: bytes) -> Dict[str, str]:
        """
        Один вызов Vision:
        - классифицируем изображение
        - извлекаем текст (если есть)
        - если текста нет — даём описание
        Возвращаем JSON-like dict.
        """
        model = self._vision_model()
        data_url = "data:image/png;base64," + base64.b64encode(img_bytes).decode("ascii")

        prompt = (
            "Проанализируй изображение.\n"
            "1) Определи тип: document / table / image / mixed.\n"
            "2) Если виден текст (в документе/таблице/скрине) — извлеки его максимально полно.\n"
            "3) Если текста нет или почти нет — дай краткое описание того, что изображено.\n"
            "Ответ верни строго в JSON (без пояснений), поля:\n"
            '{ "kind": "...", "text": "...", "description": "..." }\n'
            "Требования:\n"
            "- В поле kind используй только: document, table, image, mixed.\n"
            "- В text клади только извлечённый текст (если есть).\n"
            "- В description — 2–6 предложений, что на изображении.\n"
        )

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "input_text", "text": prompt},
                    {"type": "input_image", "image_url": data_url},
                ],
            }
        ]

        try:
            raw = self._openai.generate_text(
                model=model,
                messages=messages,
                temperature=0.0,
                max_output_tokens=1200,
                reasoning_effort=getattr(self._cfg, "openai_reasoning_effort", None),
            ) or ""
        except Exception as e:
            log.warning("Vision extract failed: %s", e)
            return {"kind": "image", "text": "", "description": ""}

        raw = raw.strip()
        # очень простой парсер: вытащим поля по ключам, чтобы не падать на “неидеальном JSON”
        def pick(key: str) -> str:
            m = re.search(rf'"{key}"\s*:\s*"((?:[^"\\]|\\.)*)"', raw, re.DOTALL)
            if not m:
                return ""
            s = m.group(1)
            s = s.replace('\\"', '"').replace("\\n", "\n").replace("\\t", "\t").replace("\\\\", "\\")
            return s

        kind = pick("kind") or "image"
        text = pick("text") or ""
        desc = pick("description") or ""
        return {"kind": kind, "text": text, "description": desc}
