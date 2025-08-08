from __future__ import annotations
import tempfile, fitz
from typing import Tuple
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

def extract_text_from_bytes(data: bytes, mime: str) -> tuple[str, dict]:
    mime = (mime or '').lower()
    if 'pdf' in mime:
        return _pdf_text(data), {"type": "pdf"}
    if 'word' in mime or 'docx' in mime:
        return _docx_text(data), {"type": "docx"}
    if 'presentation' in mime or 'pptx' in mime:
        return _pptx_text(data), {"type": "pptx"}
    if 'spreadsheet' in mime or 'excel' in mime or 'xlsx' in mime:
        return _xlsx_text(data), {"type": "xlsx"}
    try:
        s = data.decode('utf-8', errors='ignore')
        return s, {"type": "txt"}
    except:
        return "", {"type": "unknown"}

def _pdf_text(data: bytes) -> str:
    with fitz.open(stream=data, filetype='pdf') as doc:
        texts = []
        for page in doc:
            texts.append(page.get_text('text'))
        return "\n".join(texts)

def _docx_text(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".docx") as tmp:
        tmp.write(data); tmp.flush()
        doc = DocxDocument(tmp.name)
        return "\n".join(p.text for p in doc.paragraphs)

def _pptx_text(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".pptx") as tmp:
        tmp.write(data); tmp.flush()
        prs = Presentation(tmp.name)
        buf = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, 'text'):
                    buf.append(shp.text)
        return "\n".join(buf)

def _xlsx_text(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".xlsx") as tmp:
        tmp.write(data); tmp.flush()
        wb = load_workbook(tmp.name, read_only=True, data_only=True)
        buf = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) if v is not None else '' for v in row]
                if any(vals):
                    buf.append(" | ".join(vals))
        return "\n".join(buf)
